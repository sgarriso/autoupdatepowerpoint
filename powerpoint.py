from parser import Parser
from pptx import Presentation

placeholders = ["Text Placeholder 3","Text Placeholder 4", "Text Placeholder 5", "Text Placeholder 6", "Text Placeholder 7" , "Text Placeholder 8", "Text Placeholder 9", "Text Placeholder 10", "Text Placeholder 11",  "Text Placeholder 15", "Text Placeholder 16","Text Placeholder 17","Text Placeholder 30","Text Placeholder 31","Text Placeholder 32"]

def update_powerpoint():
    parser = Parser("test.yml")
    prs = Presentation("template\Presentation.pptx")
    counter = 0
    update = []
    magic = parser.get_events()
    for shape in prs.slides[0].placeholders:
        if not shape.has_text_frame:
            continue
        print(shape.text, shape.name)
        if shape.name == "Title 39":
            shape.text = parser.title
            continue
            
        if shape.name  in placeholders:
            counter = counter + 1
            update.append(shape)
        if counter == 3:
            print(update)
            counter = 0
            event = next(magic)
            print(event.field_name)
            records = [event.name, event.time, event.title]
            for shape_counter, shape in enumerate(update):
                shape.text = records[shape_counter]
            update = []
            
        
            
            
    prs.save("magic.pptx")
    
        

update_powerpoint()


    
# save_powerpoint_as_jpg()
   